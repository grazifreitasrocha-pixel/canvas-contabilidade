from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Cores
VERDE = RGBColor(0x25, 0x69, 0x5c)
DOURADO = RGBColor(0xd2, 0xae, 0x6d)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF5, 0xF5, 0xF5)
VERDE_ESCURO = RGBColor(0x1a, 0x4a, 0x42)
TEXTO_ESCURO = RGBColor(0x1e, 0x2d, 0x2a)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── HELPERS ───────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=BRANCO,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_bullet_slide(slide, items, left, top, width, height,
                     font_size=18, color=BRANCO, bullet_color=DOURADO):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        # bullet dot
        dot = p.add_run()
        dot.text = "● "
        dot.font.size = Pt(font_size - 2)
        dot.font.color.rgb = bullet_color
        dot.font.name = "Calibri"
        # text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox

def bg_solid(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_divider(slide, left, top, width, color=DOURADO, thickness=3):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, left, top, width, Pt(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

W = prs.slide_width
H = prs.slide_height
blank_layout = prs.slide_layouts[6]  # totalmente em branco

# ═══════════════════════════════════════════════════════════
# SLIDE 1 – CAPA
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg_solid(slide, VERDE)

# Bloco dourado lateral esquerdo
add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, DOURADO)

# Bloco verde escuro no rodapé
add_rect(slide, Inches(0.35), Inches(5.8), W, Inches(1.7), VERDE_ESCURO)

# Nome da empresa
add_text(slide, "CANVAS CONTABILIDADE",
         Inches(1), Inches(1.2), Inches(9), Inches(1),
         font_size=38, bold=True, color=DOURADO, align=PP_ALIGN.LEFT)

# Divisor
add_divider(slide, Inches(1), Inches(2.4), Inches(8))

# Título principal
add_text(slide, "Como economizar impostos\nde forma legal e estratégica",
         Inches(1), Inches(2.7), Inches(9), Inches(1.8),
         font_size=30, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

# Subtítulo
add_text(slide, "Apresentação para Empresários",
         Inches(1), Inches(4.6), Inches(7), Inches(0.6),
         font_size=16, color=DOURADO, align=PP_ALIGN.LEFT, italic=True)

# Rodapé
add_text(slide, "Grazielle Freitas  |  grazifreitasrocha@gmail.com",
         Inches(1), Inches(6.1), Inches(10), Inches(0.6),
         font_size=13, color=RGBColor(0xb0, 0xca, 0xc5), align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 – O PROBLEMA
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg_solid(slide, VERDE)
add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, DOURADO)

add_text(slide, "O PROBLEMA",
         Inches(0.7), Inches(0.4), Inches(10), Inches(0.6),
         font_size=13, bold=True, color=DOURADO, align=PP_ALIGN.LEFT)

add_text(slide, "Você está pagando mais\nimpostos do que deveria?",
         Inches(0.7), Inches(1.0), Inches(9.5), Inches(1.8),
         font_size=30, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

add_divider(slide, Inches(0.7), Inches(2.9), Inches(6))

add_bullet_slide(slide, [
    "O Brasil tem mais de 90 tributos diferentes",
    "Empresários pagam em média 34% do faturamento em impostos",
    "Muitos estão no regime tributário errado para seu negócio",
    "Planejamento tributário ainda é visto como privilégio de grandes empresas",
], Inches(0.7), Inches(3.1), Inches(11), Inches(3.5),
   font_size=17, color=BRANCO)

# Bloco destaque
add_rect(slide, Inches(8.5), Inches(1.0), Inches(4.2), Inches(1.8), DOURADO)
add_text(slide, "34%\ndo faturamento\nem impostos",
         Inches(8.6), Inches(1.05), Inches(4), Inches(1.7),
         font_size=20, bold=True, color=VERDE_ESCURO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 – POR QUE ISSO ACONTECE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg_solid(slide, CINZA_CLARO)
add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, VERDE)
add_rect(slide, Inches(0.35), Inches(0), W, Inches(1.1), VERDE)

add_text(slide, "POR QUE ISSO ACONTECE?",
         Inches(0.7), Inches(0.2), Inches(10), Inches(0.7),
         font_size=22, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

causas = [
    ("📌 Regime errado",      "Simples, Lucro Presumido ou Real — a escolha incorreta\npode custar milhares por ano"),
    ("📌 Falta de planejamento", "Sem análise prévia, a empresa paga mais do que\na lei exige"),
    ("📌 Contabilidade reativa", "Apenas registrar o que aconteceu não basta.\nÉ preciso estratégia antecipada"),
    ("📌 Desconhecimento",    "Existem dezenas de benefícios fiscais legais\nque a maioria dos empresários ignora"),
]

for i, (titulo, desc) in enumerate(causas):
    col = i % 2
    row = i // 2
    left = Inches(0.7 + col * 6.3)
    top = Inches(1.5 + row * 2.5)
    add_rect(slide, left, top, Inches(5.9), Inches(2.2), VERDE)
    add_text(slide, titulo, left + Inches(0.2), top + Inches(0.15),
             Inches(5.5), Inches(0.5), font_size=15, bold=True, color=DOURADO)
    add_text(slide, desc, left + Inches(0.2), top + Inches(0.65),
             Inches(5.5), Inches(1.4), font_size=13, color=BRANCO)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 – A SOLUÇÃO
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg_solid(slide, VERDE)
add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, DOURADO)

add_text(slide, "A SOLUÇÃO",
         Inches(0.7), Inches(0.4), Inches(10), Inches(0.6),
         font_size=13, bold=True, color=DOURADO)

add_text(slide, "Planejamento Tributário\nEstratégico",
         Inches(0.7), Inches(1.0), Inches(8), Inches(1.8),
         font_size=30, bold=True, color=BRANCO)

add_divider(slide, Inches(0.7), Inches(2.9), Inches(5.5))

add_text(slide,
         "Com a orientação certa, sua empresa pode reduzir\na carga tributária legalmente, aproveitando todos\nos incentivos e escolhas de regime disponíveis.",
         Inches(0.7), Inches(3.1), Inches(7.5), Inches(2),
         font_size=17, color=BRANCO)

# Caixa destaque
add_rect(slide, Inches(8.8), Inches(1.2), Inches(4), Inches(4.5), DOURADO)
add_text(slide, "Economia\nde até\n\n30%\n\nna carga\ntributária",
         Inches(9.0), Inches(1.4), Inches(3.6), Inches(4.1),
         font_size=22, bold=True, color=VERDE_ESCURO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 – COMO A CANVAS AJUDA
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg_solid(slide, CINZA_CLARO)
add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, VERDE)
add_rect(slide, Inches(0.35), Inches(0), W, Inches(1.1), VERDE)

add_text(slide, "COMO A CANVAS CONTABILIDADE AJUDA VOCÊ",
         Inches(0.7), Inches(0.2), Inches(12), Inches(0.7),
         font_size=20, bold=True, color=BRANCO)

servicos = [
    ("🎯", "Diagnóstico\nTributário",    "Análise completa do\nseu regime atual"),
    ("📊", "Escolha do\nMelhor Regime",  "Simples, Presumido\nou Real — o ideal\npara seu negócio"),
    ("💰", "Recuperação\nde Créditos",   "Identificação de\ntributos pagos\na mais"),
    ("📋", "Planejamento\nAnual",         "Estratégia fiscal\npara o ano todo\ncom antecedência"),
    ("🤝", "Consultoria\nContinuada",     "Suporte constante\npara decisões\nque impactam impostos"),
]

for i, (emoji, titulo, desc) in enumerate(servicos):
    left = Inches(0.5 + i * 2.55)
    top = Inches(1.3)
    add_rect(slide, left, top, Inches(2.35), Inches(5.5), VERDE)
    add_text(slide, emoji, left + Inches(0.1), top + Inches(0.3),
             Inches(2.1), Inches(0.6), font_size=28, align=PP_ALIGN.CENTER)
    add_divider(slide, left + Inches(0.3), top + Inches(1.05), Inches(1.6), DOURADO, 2)
    add_text(slide, titulo, left + Inches(0.1), top + Inches(1.2),
             Inches(2.1), Inches(0.9), font_size=14, bold=True,
             color=DOURADO, align=PP_ALIGN.CENTER)
    add_text(slide, desc, left + Inches(0.1), top + Inches(2.2),
             Inches(2.1), Inches(2.5), font_size=12, color=BRANCO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 – RESULTADOS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg_solid(slide, VERDE)
add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, DOURADO)

add_text(slide, "RESULTADOS REAIS",
         Inches(0.7), Inches(0.4), Inches(10), Inches(0.6),
         font_size=13, bold=True, color=DOURADO)

add_text(slide, "O que nossos clientes conquistaram",
         Inches(0.7), Inches(1.0), Inches(10), Inches(0.9),
         font_size=28, bold=True, color=BRANCO)

add_divider(slide, Inches(0.7), Inches(2.0), Inches(5))

numeros = [
    ("R$ 28 mil",   "economizados/ano\npor comércio varejista"),
    ("27%",         "de redução de\ncarga tributária"),
    ("+200",        "clientes atendidos\ncom planejamento"),
    ("100%",        "dentro da\nlegalidade"),
]

for i, (numero, desc) in enumerate(numeros):
    left = Inches(0.7 + i * 3.15)
    top = Inches(2.3)
    add_rect(slide, left, top, Inches(2.8), Inches(2.3), VERDE_ESCURO)
    add_text(slide, numero, left + Inches(0.1), top + Inches(0.2),
             Inches(2.6), Inches(0.9), font_size=26, bold=True,
             color=DOURADO, align=PP_ALIGN.CENTER)
    add_divider(slide, left + Inches(0.4), top + Inches(1.1), Inches(2.0), DOURADO, 2)
    add_text(slide, desc, left + Inches(0.1), top + Inches(1.25),
             Inches(2.6), Inches(0.9), font_size=12,
             color=BRANCO, align=PP_ALIGN.CENTER)

add_text(slide,
         '"A Canvas identificou que estávamos no regime errado há 3 anos.\nEconomizamos R$ 28 mil só no primeiro ano de ajuste." — Cliente Canvas',
         Inches(0.7), Inches(5.0), Inches(12), Inches(1.5),
         font_size=14, italic=True, color=RGBColor(0xb0, 0xca, 0xc5))

# ═══════════════════════════════════════════════════════════
# SLIDE 7 – COMO FUNCIONA
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg_solid(slide, CINZA_CLARO)
add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, VERDE)
add_rect(slide, Inches(0.35), Inches(0), W, Inches(1.1), VERDE)

add_text(slide, "COMO FUNCIONA NA PRÁTICA",
         Inches(0.7), Inches(0.2), Inches(12), Inches(0.7),
         font_size=20, bold=True, color=BRANCO)

passos = [
    ("1", "Diagnóstico\nGratuito",       "Análise sem compromisso\nda situação atual\nda sua empresa"),
    ("2", "Proposta\nPersonalizada",     "Elaboramos um plano\ntributário sob medida\npara o seu negócio"),
    ("3", "Implementação",               "Executamos todas as\nadequações e ajustes\nnecessários"),
    ("4", "Monitoramento\nContínuo",      "Acompanhamento mensal\ngarantindo os resultados\nao longo do ano"),
]

arrow_top = Inches(3.1)
for i, (num, titulo, desc) in enumerate(passos):
    left = Inches(0.6 + i * 3.17)
    top = Inches(1.4)
    # Círculo numerado (simulado com retângulo arredondado)
    add_rect(slide, left + Inches(0.8), top, Inches(0.8), Inches(0.8), DOURADO)
    add_text(slide, num, left + Inches(0.8), top,
             Inches(0.8), Inches(0.8), font_size=20, bold=True,
             color=VERDE_ESCURO, align=PP_ALIGN.CENTER)
    add_rect(slide, left, top + Inches(1.0), Inches(2.9), Inches(4.5), VERDE)
    add_text(slide, titulo, left + Inches(0.1), top + Inches(1.1),
             Inches(2.7), Inches(0.9), font_size=15, bold=True,
             color=DOURADO, align=PP_ALIGN.CENTER)
    add_divider(slide, left + Inches(0.5), top + Inches(2.1), Inches(1.9), DOURADO, 2)
    add_text(slide, desc, left + Inches(0.1), top + Inches(2.3),
             Inches(2.7), Inches(2.0), font_size=12,
             color=BRANCO, align=PP_ALIGN.CENTER)
    # seta
    if i < 3:
        add_text(slide, "▶", left + Inches(3.0), top + Inches(1.8),
                 Inches(0.4), Inches(0.5), font_size=16, color=DOURADO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 – CTA FINAL
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg_solid(slide, VERDE)
add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, DOURADO)

add_text(slide, "PRÓXIMO PASSO",
         Inches(0.7), Inches(0.5), Inches(10), Inches(0.6),
         font_size=13, bold=True, color=DOURADO)

add_text(slide, "Agende uma\nConsultoria Gratuita",
         Inches(0.7), Inches(1.2), Inches(9.5), Inches(2),
         font_size=34, bold=True, color=BRANCO)

add_divider(slide, Inches(0.7), Inches(3.3), Inches(5.5))

add_text(slide,
         "Em 30 minutos, identificamos oportunidades reais\nde economia fiscal para o seu negócio.\nSem compromisso, sem custo.",
         Inches(0.7), Inches(3.5), Inches(7.5), Inches(2),
         font_size=17, color=BRANCO)

# Caixa contato
add_rect(slide, Inches(8.5), Inches(1.2), Inches(4.5), Inches(4.5), DOURADO)
add_text(slide, "📱 WhatsApp\n\ngrazifreitasrocha\n@gmail.com\n\nCANVAS\nCONTABILIDADE",
         Inches(8.6), Inches(1.4), Inches(4.2), Inches(4.1),
         font_size=15, bold=False, color=VERDE_ESCURO, align=PP_ALIGN.CENTER)

add_text(slide, "Grazielle Freitas  —  Canvas Contabilidade",
         Inches(0.7), Inches(6.8), Inches(10), Inches(0.5),
         font_size=11, color=RGBColor(0xb0, 0xca, 0xc5))

# ─── SALVAR ────────────────────────────────────────────────
output = r"C:\Users\grazi.RAPHAEL\Downloads\Robo\Canvas_Contabilidade_Apresentacao.pptx"
prs.save(output)
print(f"Apresentação salva em: {output}")
